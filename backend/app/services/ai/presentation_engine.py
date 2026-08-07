from io import BytesIO
import os
import tempfile

import whisper
import ffmpeg

from fastapi import UploadFile

from pypdf import PdfReader
from pptx import Presentation

from app.services.ai.gemini_service import GeminiService

class PresentationEngine:

    def __init__(self):
        self.gemini = GeminiService()
        self.whisper = whisper.load_model("base")

    def speech_to_text(self, filename, content):
        with tempfile.TemporaryDirectory() as temp:
            input_file = os.path.join(temp, filename)
            with open(input_file, "wb") as f:
                f.write(content)

            audio_file = input_file
            if filename.endswith(".mp4"):
                audio_file = os.path.join(temp, "audio.wav")
                (
                    ffmpeg
                    .input(input_file)
                    .output(audio_file, ac=1, ar=16000)
                    .overwrite_output()
                    .run(quiet=True)
                )

            result = self.whisper.transcribe(audio_file)
            return result["text"]

    async def analyze(self, file: UploadFile | None = None, transcript: str | None = None):
        extracted_text = ""

        if transcript:
            extracted_text = transcript
        elif file:
            filename = file.filename.lower()
            content = await file.read()

            if filename.endswith(".pdf"):
                reader = PdfReader(BytesIO(content))
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        extracted_text += text + "\n"
            elif filename.endswith(".pptx"):
                presentation = Presentation(BytesIO(content))
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_text += shape.text + "\n"
            elif filename.endswith(".txt"):
                extracted_text = content.decode("utf-8", errors="ignore")
            elif (
                filename.endswith(".mp3")
                or filename.endswith(".wav")
                or filename.endswith(".mp4")
            ):
                extracted_text = self.speech_to_text(filename, content)
            else:
                extracted_text = filename

        if extracted_text.strip() == "":
            extracted_text = "No readable content found."

        prompt = f"""
You are an expert AI Presentation Coach.

Analyze the presentation below.

Evaluate:

- Clarity (0-100)
- Confidence (0-100)
- Speaking Speed
- Filler Words
- Strengths
- Weaknesses
- Coaching Feedback
- Overall Score

Return ONLY valid JSON.

{{
"clarity":90,
"confidence":88,
"speaking_speed":"Good",
"filler_words":["um","uh"],
"strengths":["Clear flow","Confident delivery"],
"weaknesses":["Need more evidence"],
"feedback":"Excellent presentation. Improve supporting examples.",
"overall_score":89
}}

Presentation:

{extracted_text}
"""

        return await self.gemini.generate_json(
            prompt
        )