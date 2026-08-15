"""
Real text extraction from uploaded presentation files. No content is
invented — if a slide/page has no extractable text, it's reported as empty
rather than filled in.
"""
from pptx import Presentation
import pdfplumber


def parse_pptx(file_path: str) -> list[dict]:
    """Returns [{slide_number, text}] — real text pulled from every shape on every slide."""
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        texts.append(line.strip())
        slides.append({"slide_number": i, "text": "\n".join(texts)})
    return slides


def parse_pdf(file_path: str) -> list[dict]:
    """Returns [{slide_number, text}] — one entry per PDF page, same shape as parse_pptx
    so the rest of the pipeline doesn't need to know which file type it got."""
    pages = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            pages.append({"slide_number": i, "text": text.strip()})
    return pages


def parse_document(file_path: str, filename: str) -> list[dict]:
    lower = filename.lower()
    if lower.endswith(".pptx"):
        return parse_pptx(file_path)
    elif lower.endswith(".pdf"):
        return parse_pdf(file_path)
    else:
        raise ValueError(f"Unsupported file type: {filename}. Only .pptx and .pdf are supported.")
